#!/usr/bin/env python3
"""
Extracteurs de texte - Votre MultiFormatReader adapté
"""

import os
from typing import Optional

class TextExtractor:
    """Extracteur de texte multi-format"""
    
    def __init__(self):
        print("📚 Extracteur de texte initialisé")
        
        # Formats supportés
        self.supported_formats = {
            '.pdf': self._read_pdf,
            '.docx': self._read_docx,
            '.doc': self._read_docx,
            '.txt': self._read_txt,
            '.xlsx': self._read_xlsx,
            '.xls': self._read_xlsx,
            '.pptx': self._read_pptx,
            '.ppt': self._read_pptx
        }
    
    def extract_text(self, file_path: str) -> str:
        """Extrait le texte selon le format du fichier"""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            reader_function = self.supported_formats.get(file_ext)
            
            if reader_function:
                print(f"🔄 Extraction du texte ({file_ext.upper()})...")
                text = reader_function(file_path)
                print(f"✅ Extraction terminée: {len(text)} caractères extraits")
                return text
            else:
                print(f"❌ Format non supporté: {file_ext}")
                return ""
                
        except Exception as e:
            print(f"❌ Erreur lors de l'extraction: {e}")
            return ""
    
    def _read_pdf(self, file_path: str) -> str:
        """Lit un fichier PDF"""
        try:
            import PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                full_text = ""
                
                num_pages = len(pdf_reader.pages)
                print(f"📄 PDF contient {num_pages} page(s)")
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    full_text += page_text + "\n\n"
                
                return full_text.strip()
        except ImportError:
            print("❌ PyPDF2 non installé")
            return ""
        except Exception as e:
            print(f"❌ Erreur PDF: {e}")
            return ""
    
    def _read_docx(self, file_path: str) -> str:
        """Lit un fichier Word DOCX"""
        try:
            from docx import Document
            doc = Document(file_path)
            full_text = ""
            
            print(f"📄 Document contient {len(doc.paragraphs)} paragraphes")
            
            # Extraire le texte des paragraphes
            for paragraph in doc.paragraphs:
                full_text += paragraph.text + "\n"
            
            # Extraire le texte des tableaux
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text += cell.text + " "
                full_text += "\n"
            
            return full_text.strip()
            
        except ImportError:
            print("❌ python-docx non installé")
            return ""
        except Exception as e:
            print(f"❌ Erreur DOCX: {e}")
            return self._read_txt(file_path)
    
    def _read_txt(self, file_path: str) -> str:
        """Lit un fichier texte"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
                    print(f"📄 Fichier texte lu avec encodage {encoding}")
                    return text
            except UnicodeDecodeError:
                continue
        
        print("❌ Impossible de décoder le fichier texte")
        return ""
    
    def _read_xlsx(self, file_path: str) -> str:
        """Lit un fichier Excel"""
        try:
            import pandas as pd
            excel_file = pd.ExcelFile(file_path)
            full_text = ""
            
            print(f"📄 Excel contient {len(excel_file.sheet_names)} feuille(s)")
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                for column in df.columns:
                    full_text += f"{column}: "
                    for value in df[column].dropna():
                        full_text += str(value) + " "
                    full_text += "\n"
            
            return full_text.strip()
            
        except ImportError:
            print("❌ pandas non installé")
            return ""
        except Exception as e:
            print(f"❌ Erreur Excel: {e}")
            return ""
    
    def _read_pptx(self, file_path: str) -> str:
        """Lit un fichier PowerPoint"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = ""
            
            print(f"📄 PowerPoint contient {len(prs.slides)} slide(s)")
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
            
            return full_text.strip()
            
        except ImportError:
            print("❌ python-pptx non installé")
            return ""
        except Exception as e:
            print(f"❌ Erreur PowerPoint: {e}")
            return ""